
import os
import uuid
import logging
import tempfile
import base64
from typing import TypedDict, List

import streamlit as st
from langgraph.graph import StateGraph, END

import httpx
from openai import OpenAI
from tavily import TavilyClient
import chromadb

from pptx import Presentation
from pptx.util import Inches

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

from tenacity import retry, stop_after_attempt, wait_exponential



# ---------------- LOGGING ----------------

# Structured logging for production observability
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s"
)
logger = logging.getLogger("deep_research_agent")



# ---------------- CONFIG ----------------

MAX_TOPIC_LENGTH = 500
LLM_MODEL = "gpt-5.5-mini"
EMBEDDING_MODEL = "text-embedding-3-small"
IMAGE_MODEL = "gpt-image-1"
SEARCH_MAX_RESULTS = 5
TOP_K = 3


OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
TAVILY_API_KEY = os.getenv("TAVILY_API_KEY")

# Use custom httpx client to bypass corporate proxy SSL interception
openai_client = OpenAI(api_key=OPENAI_API_KEY)
tavily_client = TavilyClient(api_key=TAVILY_API_KEY)



# ---------------- RAG STORE (ChromaDB) ----------------

# Function to initialize and cache ChromaDB collection for persistent research memory
@st.cache_resource
def get_chroma_collection():
    client = chromadb.PersistentClient(path="./chroma_db")
    
    collection = client.get_or_create_collection(
        name="research_memory",
        metadata={"hnsw:space": "cosine"}
    )
    return collection

memory_collection = get_chroma_collection()



# Function to generate embedding vector using OpenAI API

# Retries with exponential backoff for transient API failures
@retry(stop=stop_after_attempt(3), wait=wait_exponential(min=1, max=10))
def embed(text):
    """Generate embedding vector for a given text."""
    response = openai_client.embeddings.create(
        model=EMBEDDING_MODEL,
        input=text
    )
    return response.data[0].embedding



# Functions to store and retrieve research memory in ChromaDB
def store_memory(text):
    try:
        vec = embed(text)
        doc_id = str(uuid.uuid4())  # UUID avoids concurrency collisions
        memory_collection.add(
            ids=[doc_id],
            documents=[text],
            embeddings=[vec]
        )
        logger.info("Stored memory with id=%s", doc_id)

    except Exception as e:
        logger.error("Failed to store memory: %s", e)



# Function to retrieve top-K semantically similar past research from ChromaDB
def retrieve_memory(query):
    """Retrieve top-K semantically similar past research from memory."""
    try:
        if memory_collection.count() == 0:
            return ""

        q_vec = embed(query)
        results = memory_collection.query(
            query_embeddings=[q_vec],
            n_results=TOP_K
        )

        if results["documents"] and results["documents"][0]:
            return "\n\n".join(results["documents"][0])
        
    except Exception as e:
        logger.error("Failed to retrieve memory: %s", e)
    return ""




# ---------------- STATE ----------------

# LangGraph shared state — each agent reads/writes specific fields
class State(TypedDict):
    topic: str
    evidence: List[str]
    synthesis: str
    strategy: str
    refined: str
    diagram_path: str



# ---------------- LLM ----------------

@retry(stop=stop_after_attempt(3), wait=wait_exponential(min=1, max=10))
def llm(prompt):
    response = openai_client.chat.completions.create(
        model=LLM_MODEL,
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content



# ---------------- AGENTS ----------------

# Collector agent gathers evidence from Tavily web search
def collector(state):
    """Agent 1: Gather evidence from web search via Tavily."""
    topic = state["topic"]
    logger.info("Collecting evidence for: %s", topic)

    try:
        results = tavily_client.search(topic, max_results=SEARCH_MAX_RESULTS)
        evidence = [f"{r['content']} (Source: {r['url']})" for r in results["results"]]

    except Exception as e:
        logger.error("Tavily search failed: %s", e)
        evidence = [f"[Search unavailable: {e}]"]

    return {"evidence": evidence}



# Synthesizer agent combines new evidence with past memory to create structured research output
def synthesizer(state):
    """Agent 2: Combine evidence + past memory into structured research."""
    logger.info("Synthesizing research")

    try:
        memory = retrieve_memory(state["topic"])
        combined_memory = "\n".join(state["evidence"]) + "\n\nPast knowledge:\n" + memory

        prompt = f"""
            Create structured research with citations with this evidence and past knowledge: 
            {combined_memory}

            Respond in sections with clear headings. 
            Cite sources from the evidence. 
            Focus on enterprise AI architecture insights.

            Sections:
            - Product
            - Market gap
            - Architecture
            - Benefits
            - Partner implications
            - Risks
            - Future outlook
        """

        synthesis = llm(prompt)
        
    except Exception as e:
        logger.error("Synthesis failed: %s", e)
        synthesis = f"[Synthesis failed: {e}]"

    return {"synthesis": synthesis}



# Strategist agent transforms the synthesis into a consulting-style strategic report with actionable insights and roadmap
def strategist(state):
    """Agent 3: Transform research into a consulting-style strategic report."""
    logger.info("Generating strategy")
    try:
        prompt = f"""
        Convert given synthesis into a consulting report + strategic view. 
        Add roadmap, strategy, positioning, and actionable insights for enterprise decision-makers. 
        Focus on AI architecture implications.

        Synthesis:
        {state['synthesis']}
        """

        strategy = llm(prompt)
    
    except Exception as e:
        logger.error("Strategy generation failed: %s", e)
        strategy = f"[Strategy generation failed: {e}]"

    return {"strategy": strategy}



# Evaluator agent refines the report for clarity, depth, and coherence, ensuring it's polished and professional for executive consumption
def evaluator(state):
    """Agent 4: Refine report for clarity, depth, and coherence."""
    logger.info("Evaluating and refining report")
    try:
        prompt = f"""
        Improve and refine this report for clarity and depth:
        {state['strategy']}
        
        Ensure it's polished and professional for executive consumption.
        """
        refined = llm(prompt)

    except Exception as e:
        logger.error("Evaluation failed: %s", e)
        refined = state["strategy"]

    return {"refined": refined}



# Diagram Generator agent creates an architecture diagram based on the topic
def diagram_generator(state):
    """Agent 5: Generate an AI architecture diagram via image model."""
    logger.info("Generating architecture diagram")
    try:
        prompt = f"""
        Create a clean and layered architecture diagram for: {state['topic']}
        Focus on AI components, data flow, and integration points relevant to enterprise architecture.
        """

        img = openai_client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size="1024x1024",
            response_format="b64_json"
        )

        image_base64 = img.data[0].b64_json
        image_bytes = base64.b64decode(image_base64)

        # Write to temp dir to avoid CWD file collisions across sessions
        path = os.path.join(tempfile.gettempdir(), f"diagram_{uuid.uuid4().hex}.png")
        
        with open(path, "wb") as f:
            f.write(image_bytes)

        logger.info("Diagram saved to %s", path)
        return {"diagram_path": path}
    
    except Exception as e:
        logger.error("Diagram generation failed: %s", e)
    
    return {"diagram_path": ""}



# ---------------- GRAPH ----------------

# Compiled once and cached; defines the sequential agent pipeline
@st.cache_resource
def build_workflow():
    graph = StateGraph(State)
    graph.add_node("collector", collector)
    graph.add_node("synth", synthesizer)
    graph.add_node("strategy", strategist)
    graph.add_node("eval", evaluator)
    graph.add_node("diagram", diagram_generator)

    graph.set_entry_point("collector")
    graph.add_edge("collector", "synth")
    graph.add_edge("synth", "strategy")
    graph.add_edge("strategy", "eval")
    graph.add_edge("eval", "diagram")
    graph.add_edge("diagram", END)

    return graph.compile()

workflow = build_workflow()



# ---------------- EXPORT ----------------

def export_pdf(text, topic):
    """Export report as a properly formatted PDF using platypus."""
    path = os.path.join(tempfile.gettempdir(), f"report_{uuid.uuid4().hex}.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph(f"Enterprise Research: {topic}", styles["Title"]))
    story.append(Spacer(1, 20))

    for line in text.split("\n"):
        if line.strip():
            story.append(Paragraph(line, styles["BodyText"]))
            story.append(Spacer(1, 6))

    doc.build(story)
    logger.info("PDF exported to %s", path)
    return path


def export_ppt(text, diagram_path, topic):
    """Export report as a PowerPoint with content slides and diagram."""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Enterprise Research: {topic}"

    # Content slides (cap at 15 to avoid excessive slides)
    sections = [s for s in text.split("\n\n") if s.strip()]
    for sec in sections[:15]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Insight"
        s.placeholders[1].text = sec[:1000]

    # Diagram slide
    if diagram_path and os.path.exists(diagram_path):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Architecture Diagram"
        slide.shapes.add_picture(diagram_path, Inches(1), Inches(1.5), width=Inches(6))

    path = os.path.join(tempfile.gettempdir(), f"report_{uuid.uuid4().hex}.pptx")
    prs.save(path)
    logger.info("PPT exported to %s", path)
    return path


# ---------------- UI ----------------

st.title("Enterprise Research Agent")

topic = st.text_input("Enter topic", max_chars=MAX_TOPIC_LENGTH)

if st.button("Run Research") and topic:
    if len(topic.strip()) < 3:
        st.warning("Please enter a topic with at least 3 characters.")
        st.stop()

    state = {
        "topic": topic.strip(),
        "evidence": [],
        "synthesis": "",
        "strategy": "",
        "refined": "",
        "diagram_path": ""
    }

    progress_container = st.container()
    failed = False

    # Step 1: Collect evidence
    with progress_container:
        with st.spinner("🔍 Step 1/5: Collecting evidence..."):
            state.update(collector(state))
        if state["evidence"] and state["evidence"][0].startswith("[Search unavailable"):
            st.error("❌ Step 1/5: Evidence collection failed — " + state["evidence"][0])
            failed = True
        else:
            st.success("✅ Step 1/5: Evidence collected")

    # Step 2: Synthesize research
    if not failed:
        with progress_container:
            with st.spinner("🧠 Step 2/5: Synthesizing research..."):
                state.update(synthesizer(state))
            if state["synthesis"].startswith("[Synthesis failed"):
                st.error("❌ Step 2/5: Synthesis failed — " + state["synthesis"])
                failed = True
            else:
                st.success("✅ Step 2/5: Research synthesized")

    # Step 3: Generate strategy
    if not failed:
        with progress_container:
            with st.spinner("📊 Step 3/5: Generating strategy..."):
                state.update(strategist(state))
            if state["strategy"].startswith("[Strategy generation failed"):
                st.error("❌ Step 3/5: Strategy generation failed — " + state["strategy"])
                failed = True
            else:
                st.success("✅ Step 3/5: Strategy generated")

    # Step 4: Evaluate and refine
    if not failed:
        with progress_container:
            with st.spinner("✍️ Step 4/5: Evaluating and refining..."):
                state.update(evaluator(state))
            if state["refined"].startswith("["):
                st.error("❌ Step 4/5: Evaluation failed — " + state["refined"])
                failed = True
            else:
                st.success("✅ Step 4/5: Report refined")

    # Step 5: Generate diagram
    if not failed:
        with progress_container:
            with st.spinner("🎨 Step 5/5: Generating architecture diagram..."):
                state.update(diagram_generator(state))
            if not state["diagram_path"]:
                st.warning("⚠️ Step 5/5: Diagram generation failed (continuing without diagram)")
            else:
                st.success("✅ Step 5/5: Diagram generated")

    if failed:
        st.error("Pipeline stopped due to an error above. Please check your API keys and try again.")
        st.stop()

    st.divider()
    st.subheader("Final Report")
    st.write(state["refined"])

    if state["diagram_path"] and os.path.exists(state["diagram_path"]):
        st.image(state["diagram_path"], caption="Architecture Diagram")

    # Persist to memory for future topic enrichment
    store_memory(state["refined"])

    # Generate downloadable exports
    pdf = export_pdf(state["refined"], topic)
    ppt = export_ppt(state["refined"], state["diagram_path"], topic)

    col1, col2 = st.columns(2)
    with col1:
        with open(pdf, "rb") as f:
            st.download_button("📄 Download PDF", f, file_name="report.pdf", mime="application/pdf")
    with col2:
        with open(ppt, "rb") as f:
            st.download_button("📊 Download PPT", f, file_name="report.pptx",
                             mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

